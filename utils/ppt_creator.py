import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


def create_presentation(excel_path, output_path, sheet_name=None):
    # تحميل البيانات من ملف Excel
    df = pd.read_excel(excel_path, sheet_name=sheet_name)

    # إنشاء عرض تقديمي جديد
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[5]  # Layout فارغ

    for index, row in df.iterrows():
        slide = prs.slides.add_slide(title_slide_layout)
        top = Inches(1)
        left = Inches(1)
        width = Inches(8)
        height = Inches(0.5)

        # بناء النص من القيم
        content = ""
        for col in df.columns:
            content += f"{col}: {row[col]}\n"

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = content
        p.font.size = Pt(14)

    prs.save(output_path)
    print(f"✔️ Presentation saved to {output_path}")
